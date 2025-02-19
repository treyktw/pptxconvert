import os
import sys
import logging
from pathlib import Path
import win32com.client
import time
from pptx import Presentation
from typing import List, Tuple
from tqdm import tqdm
import shutil
import subprocess
import re

class PPTProcessor:
    def __init__(self, base_dir: str):
        self.base_dir = Path(base_dir)
        self.pptx_dir = self.base_dir / "pptx"
        self.text_dir = self.base_dir / "text"
        self.default_dir = self.text_dir / "default"
        self.noted_dir = self.text_dir / "noted"
        self.temp_dir = self.base_dir / ".temp"
        self.powerpoint = None
        self.setup_logging()
        self.setup_directories()

    def combine_text_files(self):
        """Combine all text files into one large file."""
        combined_file = self.text_dir / "combined_notes.txt"
        
        # First combine files from default directory
        default_files = sorted(self.default_dir.glob("*.txt"))
        noted_files = sorted(self.noted_dir.glob("*.txt"))
        
        # Prefer noted files over default files if they exist
        processed_files = set()
        
        with open(combined_file, 'w', encoding='utf-8') as outfile:
            # Write header
            outfile.write("COMBINED LECTURE NOTES\n")
            outfile.write("=" * 50 + "\n\n")
            
            # First process noted files
            for file_path in noted_files:
                base_name = file_path.stem
                processed_files.add(base_name)
                
                outfile.write(f"\nCHAPTER: {base_name}\n")
                outfile.write("-" * 50 + "\n")
                
                with open(file_path, 'r', encoding='utf-8') as infile:
                    outfile.write(infile.read())
                outfile.write("\n\n")
            
            # Then process default files that don't have notes
            for file_path in default_files:
                base_name = file_path.stem
                if base_name not in processed_files:
                    outfile.write(f"\nCHAPTER: {base_name}\n")
                    outfile.write("-" * 50 + "\n")
                    
                    with open(file_path, 'r', encoding='utf-8') as infile:
                        outfile.write(infile.read())
                    outfile.write("\n\n")
        
        self.logger.info(f"Created combined notes file: {combined_file}")
        return combined_file

    def generate_study_guide(self):
        """Generate a study guide using local LLM."""
        combined_file = self.text_dir / "combined_notes.txt"
        study_guide_file = self.base_dir / "study_guide.txt"
        
        if not combined_file.exists():
            self.logger.error("Combined notes file not found. Run combine_text_files first.")
            return
        
        try:
            # Read the combined notes
            with open(combined_file, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Create a very direct, structured prompt
            prompt = """You will create flashcards from lecture notes. Each flashcard should be on a new line in this format: Term:::Definition

For example:
Bit:::The smallest unit of digital information, representing either 0 or 1
Byte:::A unit of digital information consisting of 8 bits
Resolution:::The amount of detail in a digital image, measured in pixels

Now, create flashcards from these lecture notes:

""" + content + """

Remember: Each line must be in the format Term:::Definition with no extra text or formatting."""
            
            # Call ollama with Llama2 model and debug info
            try:
                self.logger.info("Generating study guide using Llama2...")
                result = subprocess.run(
                    ["ollama", "run", "llama2"],
                    input=prompt,
                    text=True,
                    capture_output=True,
                    encoding='utf-8'
                )
                
                if result.returncode == 0:
                    output = result.stdout.strip()
                    self.logger.info(f"Raw LLM output length: {len(output)} characters")
                    
                    # Process the output line by line
                    valid_cards = []
                    for line in output.split('\n'):
                        line = line.strip()
                        if ':::' in line and not line.startswith('```') and not line.startswith('#'):
                            term, definition = line.split(':::', 1)
                            if term.strip() and definition.strip():  # Ensure both parts exist
                                valid_cards.append(f"{term.strip()}:::{definition.strip()}")
                    
                    self.logger.info(f"Extracted {len(valid_cards)} valid flashcards")
                    
                    # If we have valid cards, write them
                    if valid_cards:
                        final_content = """DIGITAL MEDIA STUDY GUIDE
================================

Format: Term:::Definition
Ready for Quizlet Import

--------------------------------

""" + '\n'.join(valid_cards)
                        
                        with open(study_guide_file, 'w', encoding='utf-8') as f:
                            f.write(final_content)
                        self.logger.info(f"Created study guide: {study_guide_file}")
                    else:
                        # If no valid cards, try one more time with a simpler model
                        self.logger.info("No valid cards generated, trying with simpler model...")
                        return self._generate_with_backup_model()
                else:
                    raise Exception(f"LLM error: {result.stderr}")
                    
            except FileNotFoundError:
                self.logger.error("Ollama not found. Please install Ollama and the llama2 model.")
                self.logger.error("Installation instructions:")
                self.logger.error("1. Install Ollama: https://ollama.ai/")
                self.logger.error("2. Run: ollama pull llama2")
            
        except Exception as e:
            self.logger.error(f"Error generating study guide: {str(e)}")
            return self._generate_with_backup_model()
            
    def _generate_with_backup_model(self):
        """Backup method using a simpler model."""
        try:
            self.logger.info("Attempting generation with backup model (orca-mini)...")
            result = subprocess.run(
                ["ollama", "run", "orca-mini"],
                input="""Create flashcards from these lecture notes. Use exactly this format:
Term:::Definition

Example:
Bit:::The smallest unit of digital information (0 or 1)
Binary:::A number system using only 0s and 1s

Now create flashcards from these notes:

""" + content,
                text=True,
                capture_output=True,
                encoding='utf-8'
            )
            
            if result.returncode == 0 and ':::' in result.stdout:
                final_content = """DIGITAL MEDIA STUDY GUIDE
================================

Format: Term:::Definition
Ready for Quizlet Import

--------------------------------

""" + result.stdout.strip()
                
                with open(study_guide_file, 'w', encoding='utf-8') as f:
                    f.write(final_content)
                self.logger.info(f"Created study guide with backup model: {study_guide_file}")
            else:
                raise Exception("Backup model failed to generate valid content")
                
        except Exception as e:
            self.logger.error(f"Error with backup model: {str(e)}")
            with open(study_guide_file, 'w', encoding='utf-8') as f:
                f.write("Error generating study guide. Please check the log file for details.")
        self.powerpoint = None

    def setup_logging(self):
        log_file = self.base_dir / "conversion.log"
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_file, encoding='utf-8'),  # Specify UTF-8 encoding
                logging.StreamHandler()
            ]
        )
        self.logger = logging.getLogger(__name__)

    def setup_directories(self):
        """Create required directories under base directory."""
        directories = [
            self.pptx_dir,
            self.text_dir,
            self.default_dir,
            self.noted_dir,
            self.temp_dir
        ]
        
        for directory in directories:
            try:
                directory.mkdir(exist_ok=True, parents=True)
            except Exception as e:
                self.logger.error(f"Error setting up directory {directory}: {str(e)}")
                raise Exception(f"Cannot create directory: {directory}")

    def init_powerpoint(self):
        """Initialize PowerPoint application."""
        if self.powerpoint is None:
            try:
                # Kill any existing PowerPoint processes first
                os.system("taskkill /f /im powerpnt.exe 2>nul")
                time.sleep(2)  # Wait for process to fully terminate
                
                # Create new instance
                self.powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                self.powerpoint.Visible = True  # Keep window visible to avoid COM errors
                self.powerpoint.WindowState = 2  # Minimize window (1 = normal, 2 = minimized)
            except Exception as e:
                self.logger.error(f"Failed to initialize PowerPoint: {str(e)}")
                raise Exception("PowerPoint initialization failed")

    def quit_powerpoint(self):
        """Safely quit PowerPoint application and cleanup."""
        try:
            if self.powerpoint:
                try:
                    self.powerpoint.Quit()
                except:
                    pass
                finally:
                    self.powerpoint = None
            
            # Force kill any remaining PowerPoint processes
            os.system("taskkill /f /im powerpnt.exe 2>nul")
            time.sleep(2)  # Wait for process to fully terminate
        except:
            pass

    def convert_ppt_to_pptx(self, input_file: Path, retry_count=2) -> bool:
        """Convert a single PPT file to PPTX format with retries."""
        temp_output = self.temp_dir / f"{input_file.stem}_temp.pptx"
        final_output = self.pptx_dir / f"{input_file.stem}.pptx"
        
        for attempt in range(retry_count):
            try:
                self.logger.info(f"Converting to PPTX: {input_file.name} (Attempt {attempt + 1}/{retry_count})")
                
                self.init_powerpoint()  # Ensure PowerPoint is running
                
                # Open and convert
                presentation = self.powerpoint.Presentations.Open(str(input_file.absolute()))
                presentation.SaveAs(str(temp_output.absolute()), 24)
                presentation.Close()
                
                # Move to final location
                if temp_output.exists():
                    if final_output.exists():
                        final_output.unlink()
                    shutil.move(str(temp_output), str(final_output))
                
                self.logger.info(f"Successfully converted to PPTX: {input_file.name}")
                
                # Move original PPT file to processed folder or delete
                try:
                    input_file.unlink()
                    self.logger.info(f"Removed original PPT file: {input_file.name}")
                except Exception as e:
                    self.logger.error(f"Could not remove original PPT file {input_file.name}: {str(e)}")
                
                return True

            except Exception as e:
                self.logger.error(f"Error converting to PPTX {input_file.name} (Attempt {attempt + 1}): {str(e)}")
                self.quit_powerpoint()  # Reset PowerPoint instance
                time.sleep(2)  # Wait before retry
                
                # Clean up temp file if it exists
                if temp_output.exists():
                    try:
                        temp_output.unlink()
                    except:
                        pass

        return False

    def extract_text_from_shape(self, shape) -> str:
        """Extract text from a shape, handling different shape types."""
        text = ""
        try:
            if hasattr(shape, "text"):
                text = shape.text
            elif hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        row_texts.append(cell.text.strip())
                    text += " | ".join(row_texts) + "\n"
            elif hasattr(shape, "has_chart") and shape.has_chart:
                chart = shape.chart
                text = f"[Chart: {chart.chart_title if hasattr(chart, 'chart_title') else 'Untitled'}]\n"
        except Exception as e:
            self.logger.debug(f"Could not extract text from shape: {str(e)}")
        return text.strip()

    def extract_notes(self, slide) -> str:
        """Extract notes from a slide."""
        notes_text = ""
        try:
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, "notes_text_frame"):
                    notes_text = notes_slide.notes_text_frame.text
        except Exception as e:
            self.logger.debug(f"Could not extract notes: {str(e)}")
        return notes_text.strip()

    def extract_from_slide(self, slide, slide_number: int) -> Tuple[str, str]:
        """Extract all text and notes from a single slide."""
        slide_text = f"\n{'='*50}\nSlide {slide_number}\n{'='*50}\n\n"
        notes_text = ""

        for shape in slide.shapes:
            text = self.extract_text_from_shape(shape)
            if text:
                slide_text += text + "\n\n"

        notes = self.extract_notes(slide)
        if notes:
            notes_text = f"\nNotes for Slide {slide_number}:\n{notes}\n"

        return slide_text, notes_text

    def convert_pptx_to_text(self, pptx_file: Path) -> bool:
        """Convert a single PPTX file to text format."""
        temp_file = self.temp_dir / f"{pptx_file.stem}_temp.pptx"
        
        try:
            self.logger.info(f"Extracting text from: {pptx_file.name}")
            
            # Copy file to temp directory first
            shutil.copy2(pptx_file, temp_file)
            
            prs = Presentation(temp_file)
            base_name = pptx_file.stem
            main_output = self.default_dir / f"{base_name}.txt"
            notes_output = self.noted_dir / f"{base_name}.txt"
            
            main_content = []
            notes_content = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text, notes_text = self.extract_from_slide(slide, slide_number)
                main_content.append(slide_text)
                if notes_text:
                    notes_content.append(notes_text)

            # Write main content
            with open(main_output, 'w', encoding='utf-8') as f:
                f.write("\n".join(main_content))
            self.logger.info(f"Created default version: {main_output.name}")

            # Write content with notes if there are any
            if notes_content:
                full_content = main_content + ["\n\nSLIDE NOTES\n==========\n"] + notes_content
                with open(notes_output, 'w', encoding='utf-8') as f:
                    f.write("\n".join(full_content))
                self.logger.info(f"Created version with notes: {notes_output.name}")
            
            return True

        except Exception as e:
            self.logger.error(f"Error extracting text from {pptx_file.name}: {str(e)}")
            return False
            
        finally:
            if temp_file.exists():
                try:
                    temp_file.unlink()
                except:
                    pass

    def process_files(self):
        """Process all files in the base directory."""
        try:
            # Check for existing PPTX files in the base directory
            existing_pptx = list(self.base_dir.glob("*.pptx"))
            if existing_pptx:
                self.logger.info(f"Found {len(existing_pptx)} existing PPTX files")
                # Move existing PPTX files to pptx directory
                for pptx in existing_pptx:
                    target = self.pptx_dir / pptx.name
                    if not target.exists():
                        shutil.move(str(pptx), str(target))
                        self.logger.info(f"Moved existing file: {pptx.name} to pptx directory")

            # Convert PPT files if any exist
            ppt_files = list(self.base_dir.glob("*.ppt"))
            if ppt_files:
                self.logger.info(f"Found {len(ppt_files)} PPT files to convert")
                with tqdm(ppt_files, desc="Converting PPT to PPTX", unit="file") as pbar:
                    for ppt_file in pbar:
                        self.convert_ppt_to_pptx(ppt_file)

            # Process all PPTX files in the pptx directory
            pptx_files = list(self.pptx_dir.glob("*.pptx"))
            
            if not pptx_files:
                self.logger.warning("No PPTX files found to process")
                return

            self.logger.info(f"Found {len(pptx_files)} PPTX files to process")
            
            successful = 0
            failed = 0

            with tqdm(pptx_files, desc="Extracting text", unit="file") as pbar:
                for pptx_file in pbar:
                    if self.convert_pptx_to_text(pptx_file):
                        successful += 1
                    else:
                        failed += 1

            # Combine all text files into one
            self.combine_text_files()

            # Generate study guide
            self.generate_study_guide()

            self.logger.info(f"""
Processing Complete:
- Total PPTX files: {len(pptx_files)}
- Successfully processed: {successful}
- Failed: {failed}

Directory Structure:
{self.base_dir}/
|-- pptx/     (Converted PPTX files)
|-- text/
|   |-- default/  (Individual text files)
|   |-- noted/    (Text files with notes)
|   `-- combined_notes.txt  (All text combined)
`-- study_guide.txt  (Generated study guide)
""")

        finally:
            self.quit_powerpoint()  # Ensure PowerPoint is closed
            # Clean up temp directory
            try:
                shutil.rmtree(self.temp_dir)
                self.temp_dir.mkdir(exist_ok=True)
            except:
                pass

def main():
    if len(sys.argv) != 2:
        print("Usage: python script.py <directory>")
        print("Example: python script.py C:/Users/username/Documents/Chapter3")
        sys.exit(1)

    base_dir = sys.argv[1]
    processor = PPTProcessor(base_dir)
    processor.process_files()

if __name__ == "__main__":
    main()